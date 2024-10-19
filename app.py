import streamlit as st
from docx import Document
from pptx import Presentation
from transformers import pipeline
import random
import os
import re

# Function to extract text from a Word document
def extract_text_from_docx(file):
    doc = Document(file)
    full_text = []
    for para in doc.paragraphs:
        full_text.append(para.text)
    return '\n'.join(full_text)

# Function to extract text from a PowerPoint presentation
def extract_text_from_pptx(file):
    presentation = Presentation(file)
    full_text = []
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                full_text.append(shape.text)
    return '\n'.join(full_text)

# Clean and preprocess the extracted text
def preprocess_text(text):
    # Removing unnecessary white spaces, newlines, and filtering out irrelevant information
    text = re.sub(r'\n+', '\n', text.strip())  # Remove extra newlines
    text = re.sub(r'\s+', ' ', text)  # Replace multiple spaces with one
    return text

# Function to generate questions using a Hugging Face model
@st.cache_resource
def load_qg_model():
    return pipeline('text2text-generation', model='valhalla/t5-base-qa-qg-hl')

def generate_questions(text, num_questions=5):
    # Load the question-generation pipeline
    qg_model = load_qg_model()

    # Make sure `num_beams` is at least equal to `num_questions`
    num_beams = max(4, num_questions)  # Ensure beams are greater or equal to requested questions

    # Splitting the text into parts if it's too long for better processing
    text_chunks = [text[i:i+512] for i in range(0, len(text), 512)]

    # Generate questions for each chunk of text, using beam search to get multiple results
    generated_questions = []
    for chunk in text_chunks:
        generated = qg_model(
            f"generate questions: {chunk}",
            max_length=128,
            num_beams=num_beams,  # Ensure enough beams for multiple questions
            num_return_sequences=min(num_questions, num_beams)  # Return the requested number of questions
        )
        generated_questions.extend([q['generated_text'] for q in generated])
    
    return generated_questions[:num_questions]  # Return only the requested number of questions

# Function to create intelligent distractor options using a model
@st.cache_resource
def load_distractor_model():
    return pipeline('fill-mask', model='bert-base-uncased')

def generate_distractors(answer, num_options=4):
    distractor_model = load_distractor_model()
    masked_sentence = answer.replace(answer.split()[-1], '[MASK]')
    distractors = distractor_model(masked_sentence, top_k=num_options)
    options = [d['token_str'] for d in distractors if d['token_str'] != answer][:num_options-1]
    options.append(answer)  # Ensure correct answer is in the options
    random.shuffle(options)  # Shuffle to randomize the correct answer position
    return options

# Function to create multiple choice quiz with options
def create_quiz(questions, num_options=4):
    quiz = []
    for question in questions:
        # Extract the answer from the question (assuming it's present at the end or within the generated text)
        answer = question.split('?')[-1].strip()  # Placeholder logic for answer extraction
        options = generate_distractors(answer, num_options)
        quiz.append({
            "question": question.split('?')[0] + '?',
            "options": options,
            "answer": answer
        })
    return quiz

# Function to create an answer key
def create_answer_key(quiz):
    return {i+1: q['answer'] for i, q in enumerate(quiz)}

# Main app function
def main():
    st.title("AI-Based Quiz Generator")

    # File uploader for the document
    uploaded_file = st.file_uploader("Upload your document", type=["docx", "pptx"])
    
    # If the user uploads a file
    if uploaded_file is not None:
        # Determine the file type and extract text accordingly
        if uploaded_file.name.endswith(".docx"):
            text = extract_text_from_docx(uploaded_file)
        elif uploaded_file.name.endswith(".pptx"):
            text = extract_text_from_pptx(uploaded_file)

        # Preprocess the extracted text
        text = preprocess_text(text)

        # Display the extracted text
        st.write("Extracted Text:")
        st.write(text)

        # Number input for the number of quiz questions
        num_questions = st.number_input("How many questions would you like in the quiz?", min_value=1, max_value=20, step=1)

        # Generate the quiz when the user clicks the button
        if st.button("Generate Quiz"):
            questions = generate_questions(text, num_questions)
            quiz = create_quiz(questions)

            st.write("Quiz:")
            for i, q in enumerate(quiz):
                st.write(f"Q{i+1}: {q['question']}")
                for opt in q['options']:
                    st.write(f"- {opt}")

            # Generate and display the answer key
            answer_key = create_answer_key(quiz)
            st.write("Answer Key:")
            st.write(answer_key)

if __name__ == "__main__":
    main()
